from pptx import Presentation
from pptx.util import Inches as I
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Pt
import numpy as np
import pandas as pd

class Generate:
  '''
  This variable is a pandas dataframe of example data in order to run a test
  '''
  example = pd.DataFrame({'Name':['Example 1','Example 2'],
                     'Supplier':['S1','S1'],
                     'Country':['China','USA'],
                     'Tier':[1,1],
                     'Total Energy':[26500,48000],
                     'Solar Energy':[5240,7200],
                     'Wind Energy':[4265,4860],
                     'Carbon Emissions':[2480,4000],
                     'Carbon Reduced':[428,600],
                     'Factory Commitment':['Factory Example 1 Commits to having 75% of their energy consumption being renewable by 2030.','Factory Example 2 Commits to reducing carbon emissions by 80% by 2028.'],
                     'Supplier Commitment':['Supplier 1 Commits to having 75% of their energy consumption being renewable by 2030.','Supplier 1 Commits to reducing carbon emissions by 80% by 2028.'],
                     'RUI':[5.4/10,2.3/10]})
  '''
  This method takes factory data and turns it into supplier data
  @param factory_data (df) - pandas dataframe of factory data with preassumed collumns
  @return supplier_data (df) - pandas dataframe with supplier data from factory data
  '''
  def setupSupplierData(factory_data):
    supplier_data = pd.DataFrame()
    todrop = []
    for row in range(factory_data.shape[0]):
      if (factory_data['Total Energy'][row] == 0.0):
        todrop.append(row)
      elif (factory_data['Carbon Emissions'][row] == 0.0):
        todrop.append(row)
    for drop in todrop[::-1]:
      factory_data.drop(drop,axis = 0,inplace = True)
    for supplier in Generate.unique(factory_data['Supplier']):
      newrow = pd.DataFrame({'Supplier':[supplier],
                            'Country':[factory_data[factory_data['Supplier'] == supplier]['Country'].tolist()],
                            'Tier':[factory_data[factory_data['Supplier'] == supplier]['Tier'].tolist()],
                            'Total Energy':[sum(factory_data[factory_data['Supplier'] == supplier]['Total Energy'].tolist())],
                            'Solar Energy':[sum(factory_data[factory_data['Supplier'] == supplier]['Solar Energy'].tolist())],
                            'Wind Energy':[sum(factory_data[factory_data['Supplier'] == supplier]['Wind Energy'].tolist())],
                            'Carbon Emissions':[sum(factory_data[factory_data['Supplier'] == supplier]['Carbon Emissions'].tolist())],
                            'Carbon Reduced':[sum(factory_data[factory_data['Supplier'] == supplier]['Carbon Reduced'].tolist())],
                            'Supplier Commitment':[factory_data[factory_data['Supplier'] == supplier]['Supplier Commitment'].tolist()[0]],
                            'RUI':[sum(factory_data[factory_data['Supplier'] == supplier]['RUI'].tolist())]})
      supplier_data = pd.concat([supplier_data,newrow])
    return supplier_data

  '''
  This method takes an entry of a string and turns it into paragraph form 
  where the longest length of a row is 48 characters
  @param: entry (str) - String input to be turned into paragraph form.
  @return: toreturn (str) - Converted string with newline's added and, if too long
                            truncated with ... at the end.
  '''
  def returnParagraph(entry:str):
    toreturn = ""
    start = 0
    for i in range(int(len(entry)/48)+1):
      toreturn += entry[start:(48*(i+1))] + '\n'
      start = (48*(i+1))
      if i == 4:
        toreturn = toreturn[:-2] + '...'
        break
    return toreturn

  '''
  This method turns an RUI score into an RUI level.
  @param RUIscore (float) - A float number representative of the RUI score
  @return toreturn (int) - The level (1, 2, or 3) the score represents
  '''
  def RUILevel(RUIscore:float):
    toreturn = int((((RUIscore)*10)/3)+1)
    return toreturn

  '''
  This method find all unique values in a list
  @param list1 (list) - list of strings
  @return unique_list (list) - list of unique values
  '''
  def unique(list1):
      list_set = set(list1)
      unique_list = (list(list_set))
      return unique_list

  '''
  This method determines the percentile in which a number falls in a list of numbers
  @param lst (list) - a list of numbers
  @param value (int) - the number in the list to find the percentile of
  @return toreturn (int) - the percentile the number falls into
  '''
  def percentile(lst,value):
    toreturn = round(lst.index(value)/len(lst)*100,0)
    return toreturn
  '''
  This method creates dataframes for the supplier templates
  @param dfrow (df.iloc) - dataframe location from supplier_data
  @param supplier_data (df) - supplier data
  @return Supplier_Template (df) - customized supplier template for the given row
  '''
  def createTemplateSupplier(dfrow,supplier_data):
    t,l,c,fl = 'Text','Line','Circle','Fill Line'
    Supplier_Template = pd.DataFrame({'Type':[l,t,t,t,l,t,t,t,t,l,l,t,t,
                                              t,l,l,l,c,c,t,t,t,t,t,
                                              t,l,t,t,
                                              t,t,fl,l,t,
                                              t,t,t,t],
                                    'Left':[5.5,9.33,.33,.5,.5,.75,3.625,.75,3,1,1,.75,.75,
                                            3,1,1,.5,.625,.625,.6375,.625,1.125,1.125,2.875,
                                            2.875,.5,.75,.75,
                                            6,9.5,6.125,6.125,6.125,
                                            6,6.25,6.25,6.25],
                                    'Top':[.25,8.25,.25,.5,1.5,1.75,1.625,2.25,2.25,2.8333,2.8333,3.0833,3.5833,
                                          3.5833,4.166,4.166,4.41666,4.75,5.5,4.75,5.5,4.8125,5.5625,4.8125,
                                          5.5625,6.25,6.5,6.95833,
                                          .5,.5,1,1,1.5,
                                          2.25,3,4.75,7],
                                    'Width':[.005,.75,.75,4.5,4.5,3.25,3.25,2.25,2.25,3.5,3.5*(dfrow['Solar Energy']+dfrow['Wind Energy'])/(dfrow['Total Energy']),3.25,2.25,
                                            2,3.5,3.5*(dfrow['Carbon Reduced'])/(dfrow['Carbon Emissions']), 4.5,.5,.5,.333,.333,1.75,1.75,2,
                                            2,4.5,3.5,3.5,
                                            3,1,4,4*(dfrow['RUI']),1,
                                            1,1,1,1],
                                    'Height':[8,.75,.75,.75,.005,.5,.5,.333,.333,.005,.05,.5,.333,
                                              .333,.005,.05,.005,.5,.5,.333,.333,.25,.25,.25,
                                              .25,.005,.333,.333,
                                              .5,.5,.25,.25,.5,
                                              .5,.5,.5,1],
                                    'Text':[None,'Supplier Card Information','Supplier Card Information',f"{dfrow['Supplier'][:26]}",None,"Energy Consumption",f"{round((dfrow['Solar Energy']+dfrow['Wind Energy'])/(dfrow['Total Energy'])*100,0)}%",f"Renewable: {round(int(dfrow['Wind Energy'])+int(dfrow['Solar Energy']),0)} MWh",f"Total: {int(dfrow['Total Energy'])} MWh",None,None,'Carbon Emissions',f"Reduced: {round(int(dfrow['Carbon Reduced']),0)} MT/y",
                                            f"Total: {round(int(dfrow['Carbon Emissions']),0)} MT/y",None,None,None,None,None,'S',"W",f"Usage: {round(int(dfrow['Solar Energy']),0)} MWh",f"Usage: {round(int(dfrow['Wind Energy']),0)} MWh",f"% of Renewable: {round(int(dfrow['Solar Energy'])/(int(dfrow['Wind Energy']+1)+int(dfrow['Solar Energy']))*100,0)}%",
                                            f"% of Renewable: {round(int(dfrow['Wind Energy'])/(int(dfrow['Wind Energy']+1)+int(dfrow['Solar Energy']))*100,0)}%", None,'Energy Commitment:',f"{Generate.returnParagraph(dfrow['Supplier Commitment'])}",
                                            'Renewables Understanding Index:',f"Level {Generate.RUILevel(dfrow['RUI'])}",None,None,f"This supplier is in the top {Generate.percentile(supplier_data['RUI'].tolist(),dfrow['RUI'])}% of suppliers.",
                                            'Factories:',"Factories Countries (top 4):"+ "".join(['\n\t-'+str(i) for i in pd.DataFrame({1:dfrow['Country']})[1].value_counts()[:4].keys().tolist()]), f"Tiers:\n\tTier 1 - {dfrow['Tier'].count(1)} Factories\n\tTier 2 - {dfrow['Tier'].count(2)} Factories\n\tTier 3+ - {len(dfrow['Tier'])-(dfrow['Tier'].count(1) + dfrow['Tier'].count(2))} Factories","For Specific Information About A\nFactory, Refer to the Factory Card."],
                                    'Font Size':[None,10,10,30,None,20,28,16,16,None,None,20,16,
                                                16,None,None,None,None,None,24,24,14,14,14,
                                                14,None,18,14,
                                                15.25,15.25,None,None,16,
                                                20,18,18,14],
                                    'Bold':[None,None,None,None,None,None,True,None,None,None,None,None,None,
                                            None,None,None,None,None,None,None,None,None,None,None,
                                            None,None,None,None,
                                            None,None,None,None,None,
                                            None,None,None,None]})
    return Supplier_Template

  '''
  This method creates dataframes for the factory templates
  @param dfrow (df.iloc) - dataframe location from factory_data
  @param factory_data (df) - factory data
  @return Factory_Template (df) - customized factory template for the given row
  '''
  def createTemplateFactory(dfrow,factory_data):
    t,l,c,fl = 'Text','Line','Circle','Fill Line'
    Factory_Template = pd.DataFrame({'Type':[l,t,t,t,l,t,t,t,t,l,l,t,t,
                                              t,l,l,l,c,c,t,t,t,t,t,
                                              t,l,t,t,
                                              t,t,fl,l,t],
                                    'Left':[5.5,9.33,.33,.5,.5,.75,3.625,.75,3,1,1,.75,.75,
                                            3,1,1,.5,.625,.625,.6375,.625,1.125,1.125,2.875,
                                            2.875,.5,.75,.75,
                                            6,9.5,6.125,6.125,6.125],
                                    'Top':[.25,8.25,.25,.5,1.5,1.75,1.625,2.25,2.25,2.8333,2.8333,3.0833,3.5833,
                                          3.5833,4.166,4.166,4.41666,4.75,5.5,4.75,5.5,4.8125,5.5625,4.8125,
                                          5.5625,6.25,6.5,6.95833,
                                          .5,.5,1,1,1.5],
                                    'Width':[.005,.75,.75,4.5,4.5,3.25,3.25,2.25,2.25,3.5,3.5*(dfrow['Solar Energy']+dfrow['Wind Energy'])/(dfrow['Total Energy']),3.25,2.25,
                                            2,3.5,3.5*(dfrow['Carbon Reduced'])/(dfrow['Carbon Emissions']), 4.5,.5,.5,.333,.333,1.75,1.75,2,
                                            2,4.5,3.5,3.5,
                                            3,1,4,4*(dfrow['RUI']),1],
                                    'Height':[8,.75,.75,.75,.005,.5,.5,.333,.333,.005,.05,.5,.333,
                                              .333,.005,.05,.005,.5,.5,.333,.333,.25,.25,.25,
                                              .25,.005,.333,.333,
                                              .5,.5,.25,.25,.5],
                                    'Text':[None,'Factory Card Information','Factory Card Information',f"{dfrow['Name'][:26]}",None,"Energy Consumption",f"{round((dfrow['Solar Energy']+dfrow['Wind Energy'])/(dfrow['Total Energy'])*100,0)}%",f"Renewable: {round(int(dfrow['Wind Energy'])+int(dfrow['Solar Energy']),0)} MWh",f"Total: {int(dfrow['Total Energy'])} MWh",None,None,'Carbon Emissions',f"Reduced: {round(int(dfrow['Carbon Reduced']),0)} MT/y",
                                            f"Total: {round(int(dfrow['Carbon Emissions']),0)} MT/y",None,None,None,None,None,'S',"W",f"Usage: {round(int(dfrow['Solar Energy']),0)} MWh",f"Usage: {round(int(dfrow['Wind Energy']),0)} MWh",f"% of Renewable: {round(int(dfrow['Solar Energy'])/(int(dfrow['Wind Energy']+1)+int(dfrow['Solar Energy']))*100,0)}%",
                                            f"% of Renewable: {round(int(dfrow['Wind Energy'])/(int(dfrow['Wind Energy']+1)+int(dfrow['Solar Energy']))*100,0)}%", None,'Energy Commitment:',f"{Generate.returnParagraph(dfrow['Factory Commitment'])}",
                                            'Renewables Understanding Index:',f"Level {Generate.RUILevel(dfrow['RUI'])}",None,None,f"This factory is in the top {Generate.percentile(factory_data['RUI'].tolist(),dfrow['RUI'])}% of factories."],
                                    'Font Size':[None,10,10,30,None,20,28,16,16,None,None,20,16,
                                                16,None,None,None,None,None,24,24,14,14,14,
                                                14,None,18,14,
                                                15.25,15.25,None,None,16],
                                    'Bold':[None,None,None,None,None,None,True,None,None,None,None,None,None,
                                            None,None,None,None,None,None,None,None,None,None,None,
                                            None,None,None,None,
                                            None,None,None,None,None]})
    return Factory_Template

  '''
  This method is used to create a text box on a presentation slide with given factors
  @param slide (presentaiton) - the current slide the information is going on 
  @param dfrow (df.iloc) - given supplier or factory data
  @param left (float) - how far from the left side of the slide the text will start
  @param top (float) - how far from the top of the slide the text will start
  @param width (float) - the width of the text box
  @param height (float) - the height of the text box
  @param text (str) - A string with the words shown within the text box
  @param fontsize (int) - the font size
  @param bold (Boolean/None) - whether the text is bold. If not bold use None type, if bold use True. (default - None)
  @param fonttype (str) - The name of the font type (default - 'Calibri')
  @return slide (presentaiton) - updated version of the slide enterd as a parameter
  '''
  def text(slide,dfrow,left,top,width,height,text,fontsize,bold = None ,fonttype = 'Calibri'):
    txBox = slide.shapes.add_textbox(I(left),I(top),I(width),I(height))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = fonttype
    font.size = Pt(fontsize)
    font.bold = bold
    return slide
  '''
  This method is used to create a line on a presentation slide with given factors
  @param slide (presentaiton) - the current slide the information is going on 
  @param dfrow (df.iloc) - given supplier or factory data
  @param left (float) - how far from the left side of the slide the line will start
  @param top (float) - how far from the top of the slide the line will start
  @param width (float) - the width of the line
  @param height (float) - the height of the line
  @return slide (presentaiton) - updated version of the slide enterd as a parameter
  '''
  def line(slide,dfrow,left,top,width,height):
    line = slide.shapes.add_shape(MSO_CONNECTOR.STRAIGHT, I(left),I(top),I(width),I(height))
    return slide

  '''
  This method is used to create a filled line on a presentation slide with given factors
  @param slide (presentaiton) - the current slide the information is going on 
  @param dfrow (df.iloc) - given supplier or factory data
  @param left (float) - how far from the left side of the slide the line will start
  @param top (float) - how far from the top of the slide the line will start
  @param width (float) - the width of the line
  @param height (float) - the height of the line
  @return slide (presentaiton) - updated version of the slide enterd as a parameter
  '''
  def fillLine(slide,dfrow,left,top,width,height):
    line = slide.shapes.add_shape(MSO_CONNECTOR.STRAIGHT,I(left),I(top),I(width),I(height))
    fill = line.fill
    fill.background()
    return slide

  '''
  This method is used to create a circle on a presentation slide with given factors
  @param slide (presentaiton) - the current slide the information is going on 
  @param dfrow (df.iloc) - given supplier or factory data
  @param left (float) - how far from the left side of the slide the circle will start
  @param top (float) - how far from the top of the slide the circle will start
  @param width (float) - the width of the circle
  @param height (float) - the height of the circle
  @return slide (presentaiton) - updated version of the slide enterd as a parameter
  '''
  def circle(slide,dfrow,left,top,width,height):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,I(left),I(top),I(width),I(height))
    return slide

  '''
  This method utilizes the created template to add the content to the slide in order
  @param prs (presentaiton) - the full presentaiton object
  @param slide (presentaiton) - the current slide the information is going on 
  @param Sup_or_Fac (str) - whether it will be with the supplier or factory template
  @param dfrow (df.iloc) - given supplier or factory data
  @return prs (presentaiton) - the updated presentation
  @return slide (presentaiton) - the updated slide
  '''
  def run(prs,slide,Sup_or_Fac,dfrow,supplier_data,factory_data):   
    if Sup_or_Fac == 'Supplier':
      Supplier_Template = Generate.createTemplateSupplier(dfrow,supplier_data)
      for row in range(Supplier_Template.shape[0]):
        if Supplier_Template['Type'][row] == 'Text':
          slide = Generate.text(slide,dfrow,Supplier_Template['Left'][row],Supplier_Template['Top'][row],Supplier_Template['Width'][row],Supplier_Template['Height'][row],Supplier_Template['Text'][row],Supplier_Template['Font Size'][row],Supplier_Template['Bold'][row])
        elif Supplier_Template['Type'][row] == 'Line':
          slide = Generate.line(slide,dfrow,Supplier_Template['Left'][row],Supplier_Template['Top'][row],Supplier_Template['Width'][row],Supplier_Template['Height'][row])
        elif Supplier_Template['Type'][row] == 'Fill Line':
          slide = Generate.fillLine(slide,dfrow,Supplier_Template['Left'][row],Supplier_Template['Top'][row],Supplier_Template['Width'][row],Supplier_Template['Height'][row])
        elif Supplier_Template['Type'][row] == 'Circle':
          slide = Generate.circle(slide,dfrow,Supplier_Template['Left'][row],Supplier_Template['Top'][row],Supplier_Template['Width'][row],Supplier_Template['Height'][row])
    elif Sup_or_Fac == 'Factory':
      Factory_Template = Generate.createTemplateFactory(dfrow,factory_data)
      for row in range(Factory_Template.shape[0]):
        if Factory_Template['Type'][row] == 'Text':
          slide = Generate.text(slide,dfrow,Factory_Template['Left'][row],Factory_Template['Top'][row],Factory_Template['Width'][row],Factory_Template['Height'][row],Factory_Template['Text'][row],Factory_Template['Font Size'][row],Factory_Template['Bold'][row])
        elif Factory_Template['Type'][row] == 'Line':
          slide = Generate.line(slide,dfrow,Factory_Template['Left'][row],Factory_Template['Top'][row],Factory_Template['Width'][row],Factory_Template['Height'][row])
        elif Factory_Template['Type'][row] == 'Fill Line':
          slide = Generate.fillLine(slide,dfrow,Factory_Template['Left'][row],Factory_Template['Top'][row],Factory_Template['Width'][row],Factory_Template['Height'][row])
        elif Factory_Template['Type'][row] == 'Circle':
          slide = Generate.circle(slide,dfrow,Factory_Template['Left'][row],Factory_Template['Top'][row],Factory_Template['Width'][row],Factory_Template['Height'][row])
    return prs, slide

  '''
  This method utilizes other methods in order to generate a presentation with
  the designated templates for the given suppliers and factories
  @param factory_data (df) - dataframe of factory data
  @param supplier_data (df) - dataframe of supplier data
  @return prs (presentaiton) - this is the final formatted presentation
  '''
  def Generate_Presentation(factory_data,supplier_data):
    prs = Presentation()
    prs.slide_width = I(11)
    prs.slide_height = I(8.5)
    lyt=prs.slide_layouts[6]
    unq = Generate.unique(factory_data['Supplier'])
    for supplier in range(len(unq)):
      slide=prs.slides.add_slide(lyt)
      prs,slide = Generate.run(prs,slide,'Supplier',supplier_data.iloc[supplier],supplier_data,factory_data)
      factories = factory_data[factory_data['Supplier'] == unq[supplier]]
      for factory in range(len(factories['Name'])):
        slide=prs.slides.add_slide(lyt)
        prs,slide = Generate.run(prs,slide,'Factory',factory_data.iloc[factory],supplier_data,factory_data)
    return prs

  '''
  This method saves the presentaiton to the given path after it is generated
  @param prs (presentation) - the presentation fully formatted
  @param path (str) - the filepath for the presentation to be saved in
  '''
  def save_prs(prs,path):
    prs.save(path)
