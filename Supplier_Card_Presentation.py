from pptx import Presentation
from pptx.util import Inches as I
import pandas as pd
from Supplier_Card_Methods import Generate

example = pd.DataFrame({'Name':['Example 1','Example 2'],
                     'Supplier':['S1','S1'],
                     'Country':['China','USA'],
                     'Tier':[1,1],
                     'Total Energy':[26500,48000],
                     'Solar Energy':[5240,7200],
                     'Wind Energy':[4265,4860],
                     'Carbon Emissions':[2480,4000],
                     'Carbon Reduced':[428,600],
                     'Factory Commitment':['Factory Example 1 Commits to having 75% of their energy consumption being renewable by 2030.','Factory Example 2 Commits to reducing carbon emissions by 80% by 2028.'],
                     'Supplier Commitment':['Supplier 1 Commits to having 75% of their energy consumption being renewable by 2030.','Supplier 1 Commits to reducing carbon emissions by 80% by 2028.'],
                     'RUI':[5.4/10,2.3/10]})


prs = Presentation()
prs.slide_width = I(11)
prs.slide_height = I(8.5)
lyt=prs.slide_layouts[6]

factory_data = example

supplier_data = Generate.setupSupplierData(factory_data)

unq = Generate.unique(factory_data['Supplier'])
for supplier in range(len(unq)):
  slide=prs.slides.add_slide(lyt)
  prs,slide = Generate.run(prs,slide,'Supplier',supplier_data.iloc[supplier],supplier_data,factory_data)
  factories = factory_data[factory_data['Supplier'] == unq[supplier]]
  for factory in range(len(factories['Name'])):
    slide=prs.slides.add_slide(lyt)
    prs,slide = Generate.run(prs,slide,'Factory',factory_data.iloc[factory],supplier_data,factory_data)


saveTo = input("Where would you like to save the presentation to?")
prs.save(saveTo+'Supplier_and_Factory_Presentation.pptx')
